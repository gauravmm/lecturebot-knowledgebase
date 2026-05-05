"""Marshal lecturer corpus into per-lecture JSONL chunks.

For each lecture in data/raw/lecturer/lectures/<year>-<slug>/:
  - emit a chunk for the verbatim page_text from manifest.json
  - if slides.pptx exists, emit one chunk per slide body and per slide notes
  - if slides.pdf exists, emit one chunk per PDF page
  - astar-workshop ships multiple PDF decks; handle all *.pdf alongside the
    declared one

Output: data/processed/lecturer/lectures/<slug>.chunks.jsonl
"""

from __future__ import annotations

import json
from pathlib import Path

import fitz
from pptx import Presentation

RAW = Path("data/raw/lecturer/lectures")
OUT = Path("data/processed/lecturer/lectures")


def extract_pptx(path: Path):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        body_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    body_lines.append(line)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        yield i, "\n".join(body_lines), notes


def extract_pdf(path: Path):
    doc = fitz.open(path)
    try:
        for i, page in enumerate(doc, 1):
            yield i, page.get_text()
    finally:
        doc.close()


def process_lecture(manifest_path: Path) -> tuple[Path, int]:
    manifest = json.loads(manifest_path.read_text())
    slug = manifest["slug"]
    out_path = OUT / f"{slug}.chunks.jsonl"
    out_path.parent.mkdir(parents=True, exist_ok=True)

    base = {
        "corpus": "lecturer/lectures",
        "doc_title": manifest["title"],
        "source_url": manifest["url"],
        "fetched_at": manifest["fetched_at"],
        "speaker": manifest.get("speaker"),
        "venue": manifest.get("venue"),
        "course": manifest.get("course"),
        "semester": manifest.get("semester"),
    }

    chunks: list[dict] = []

    page_text = (manifest.get("page_text") or "").strip()
    if page_text:
        chunks.append({
            **base,
            "id": f"{slug}::page",
            "section_path": "page",
            "text": page_text,
        })

    raw_dir = manifest_path.parent
    # Process every PDF or PPTX in the directory (astar-workshop has multiple).
    for slide_file in sorted(raw_dir.iterdir()):
        suffix = slide_file.suffix.lower()
        if suffix == ".pptx":
            for n, body, notes in extract_pptx(slide_file):
                deck = slide_file.stem
                if body.strip():
                    chunks.append({
                        **base,
                        "id": f"{slug}::{deck}::slide-{n:03d}",
                        "section_path": f"deck:{deck}/slide:{n}",
                        "deck": deck,
                        "slide_number": n,
                        "text": body,
                    })
                if notes.strip():
                    chunks.append({
                        **base,
                        "id": f"{slug}::{deck}::slide-{n:03d}-notes",
                        "section_path": f"deck:{deck}/slide:{n}/notes",
                        "deck": deck,
                        "slide_number": n,
                        "speaker_role": "speaker_notes",
                        "text": notes,
                    })
        elif suffix == ".pdf":
            for n, text in extract_pdf(slide_file):
                if text.strip():
                    deck = slide_file.stem
                    chunks.append({
                        **base,
                        "id": f"{slug}::{deck}::page-{n:03d}",
                        "section_path": f"deck:{deck}/page:{n}",
                        "deck": deck,
                        "slide_number": n,
                        "text": text,
                    })

    with out_path.open("w") as f:
        for c in chunks:
            f.write(json.dumps(c, ensure_ascii=False) + "\n")
    return out_path, len(chunks)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    total = 0
    for manifest in sorted(RAW.glob("*/manifest.json")):
        out_path, n = process_lecture(manifest)
        total += n
        print(f"{out_path}: {n} chunks")
    print(f"\nTotal: {total} chunks across {len(list(RAW.glob('*/manifest.json')))} lectures")


if __name__ == "__main__":
    main()
